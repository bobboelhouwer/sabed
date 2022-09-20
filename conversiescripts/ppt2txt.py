from pptx import Presentation
import sys
import os
import re

startdir=sys.argv[1]

def process_files(files,root):
    for f in files:
        print('Processsing', root+f,file=sys.stderr)
        if re.search('pptx$',f):
            outfilename=f+'.txt'
            outfile = open(root+'/'+outfilename,'w') 
            try:
                prs = Presentation(root+'/'+f)
            except:
                continue
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        print(shape.text,file=outfile)


for root,subFolders,files in os.walk(startdir):
    process_files(files,root)
    
